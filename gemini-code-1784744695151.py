import streamlit as st
import PyPDF2
from google import genai
from pptx import Presentation
from pptx.util import Pt
import json
import io
import os

# Initialize Gemini Client using new Google SDK
api_key = os.environ.get("GEMINI_API_KEY")
client = genai.Client(api_key=api_key) if api_key else None

def extract_text_from_pdfs(pdf_files):
    text = ""
    for pdf in pdf_files:
        reader = PyPDF2.PdfReader(pdf)
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
    return text

def generate_insights(text):
    prompt = f"""
    Analyze the following text extracted from multiple documents. 
    1. Identify common themes, patterns, and key data points.
    2. Generate a highly concise briefing note and a list of action items.
    Use simple, clear, and direct language.
    
    Text: {text[:75000]}
    """
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt
    )
    return response.text

def generate_slide_content(text):
    prompt = f"""
    Based on the following text, create an outline for a 15-slide presentation.
    Return ONLY a JSON object with a "slides" key containing an array of 15 objects. 
    Each object must have a "title" and a "content" (a list of 3-4 short bullet points).

    Text: {text[:75000]}
    """
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config={'response_mime_type': 'application/json'}
    )
    return json.loads(response.text)

def create_ppt(slide_data):
    prs = Presentation()
    slides = slide_data.get('slides', slide_data)
    for slide_info in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get('title', 'Slide')
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        for point in slide_info.get('content', []):
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

st.set_page_config(page_title="PDF Analyzer & PPT Generator", layout="wide")
st.title("📄 Cross-PDF Analyzer & Slide Generator")

uploaded_files = st.file_uploader("Upload 5-6 PDF files", type="pdf", accept_multiple_files=True)

if uploaded_files:
    if st.button("Analyze & Generate"):
        if not api_key:
            st.error("GEMINI_API_KEY is missing! Please set it in Streamlit Secrets.")
        else:
            with st.spinner("Extracting text from PDFs..."):
                extracted_text = extract_text_from_pdfs(uploaded_files)
            
            with st.spinner("Analyzing themes and generating briefing note..."):
                try:
                    insights = generate_insights(extracted_text)
                    st.subheader("Briefing Note & Action Items")
                    st.write(insights)
                except Exception as e:
                    st.error(f"Error generating insights: {e}")
                    st.stop()
                
            with st.spinner("Generating 15-page Presentation..."):
                try:
                    slide_json = generate_slide_content(extracted_text)
                    ppt_file = create_ppt(slide_json)
                    
                    st.success("Analysis and Presentation Generation Complete!")
                    st.download_button(
                        label="📥 Download 15-Slide Presentation (.pptx)",
                        data=ppt_file,
                        file_name="Executive_Summary.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.error(f"An error occurred while generating the presentation: {e}")
